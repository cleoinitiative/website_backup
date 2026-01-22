#!/usr/bin/env python3
"""
Internal Resources Converter for CLEO RAG Corpus

Converts CLEO's internal resources (docx, pptx, pdf) to markdown
for ingestion into the RAG chatbot.
"""

import os
import re
from pathlib import Path
from typing import Optional

# Try to import document parsers
try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False
    print("Warning: python-docx not installed. Run: pip install python-docx")

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    print("Warning: python-pptx not installed. Run: pip install python-pptx")

try:
    import fitz  # PyMuPDF
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False
    print("Warning: PyMuPDF not installed. Run: pip install pymupdf")


# Configuration
RESOURCES_DIR = Path("/Users/aaronsmolyar/Documents/Resources")
OUTPUT_DIR = Path(__file__).parent.parent / "corpus" / "cleo-internal"


def clean_text(text: str) -> str:
    """Clean up extracted text."""
    # Remove excessive whitespace
    text = re.sub(r'\n{3,}', '\n\n', text)
    text = re.sub(r' {2,}', ' ', text)
    text = re.sub(r'\t+', ' ', text)
    return text.strip()


def extract_docx(file_path: Path) -> Optional[str]:
    """Extract text from a Word document."""
    if not DOCX_AVAILABLE:
        return None
    
    try:
        doc = Document(file_path)
        paragraphs = []
        
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                # Check if it's a heading
                if para.style.name.startswith('Heading'):
                    level = int(para.style.name[-1]) if para.style.name[-1].isdigit() else 1
                    paragraphs.append(f"{'#' * level} {text}")
                else:
                    paragraphs.append(text)
        
        # Also extract from tables
        for table in doc.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if cells:
                    paragraphs.append(" | ".join(cells))
        
        return clean_text('\n\n'.join(paragraphs))
    except Exception as e:
        print(f"   Error reading {file_path.name}: {e}")
        return None


def extract_pptx(file_path: Path) -> Optional[str]:
    """Extract text from a PowerPoint presentation."""
    if not PPTX_AVAILABLE:
        return None
    
    try:
        prs = Presentation(file_path)
        content = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = []
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            
            if slide_text:
                content.append(f"## Slide {slide_num}\n\n" + "\n\n".join(slide_text))
        
        return clean_text('\n\n'.join(content))
    except Exception as e:
        print(f"   Error reading {file_path.name}: {e}")
        return None


def extract_pdf(file_path: Path) -> Optional[str]:
    """Extract text from a PDF document."""
    if not PDF_AVAILABLE:
        return None
    
    try:
        doc = fitz.open(file_path)
        content = []
        
        for page_num, page in enumerate(doc, 1):
            text = page.get_text().strip()
            if text:
                content.append(f"## Page {page_num}\n\n{text}")
        
        doc.close()
        return clean_text('\n\n'.join(content))
    except Exception as e:
        print(f"   Error reading {file_path.name}: {e}")
        return None


def process_file(file_path: Path, output_dir: Path) -> bool:
    """Process a single file and save as markdown."""
    suffix = file_path.suffix.lower()
    
    # Extract content based on file type
    if suffix == '.docx':
        content = extract_docx(file_path)
    elif suffix == '.pptx':
        content = extract_pptx(file_path)
    elif suffix == '.pdf':
        content = extract_pdf(file_path)
    else:
        return False
    
    if not content or len(content) < 50:
        return False
    
    # Create output filename
    safe_name = re.sub(r'[^\w\s-]', '', file_path.stem)
    safe_name = re.sub(r'\s+', '-', safe_name).lower()[:60]
    
    # Determine category from path
    rel_path = file_path.relative_to(RESOURCES_DIR)
    category = rel_path.parts[0] if len(rel_path.parts) > 1 else "general"
    category = re.sub(r'[^\w\s-]', '', category)
    category = re.sub(r'\s+', '-', category).lower()
    
    category_dir = output_dir / category
    category_dir.mkdir(parents=True, exist_ok=True)
    
    output_path = category_dir / f"{safe_name}.md"
    
    # Handle duplicates
    counter = 1
    while output_path.exists():
        output_path = category_dir / f"{safe_name}-{counter}.md"
        counter += 1
    
    # Write markdown file
    title = file_path.stem.replace('_', ' ').replace('-', ' ').title()
    
    md_content = f"# {title}\n\n"
    md_content += f"> Source: CLEO Internal Resources - {rel_path}\n\n"
    md_content += content
    
    with open(output_path, 'w', encoding='utf-8') as f:
        f.write(md_content)
    
    return True


def main():
    """Process all internal resources."""
    print("🚀 Processing CLEO Internal Resources...")
    print(f"   Input directory: {RESOURCES_DIR}")
    print(f"   Output directory: {OUTPUT_DIR}")
    
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    
    if not RESOURCES_DIR.exists():
        print(f"❌ Resources directory not found: {RESOURCES_DIR}")
        return {"error": "Directory not found"}
    
    # Find all supported files
    files = []
    for ext in ['*.docx', '*.pptx', '*.pdf']:
        files.extend(RESOURCES_DIR.rglob(ext))
    
    # Filter out temp files
    files = [f for f in files if not f.name.startswith('~$')]
    
    print(f"   Found {len(files)} documents to process")
    
    processed = 0
    failed = 0
    
    for file_path in files:
        try:
            if process_file(file_path, OUTPUT_DIR):
                processed += 1
                print(f"   ✅ {file_path.name}")
            else:
                failed += 1
        except Exception as e:
            failed += 1
            print(f"   ❌ {file_path.name}: {e}")
    
    print(f"\n✨ Processing complete!")
    print(f"   Files processed: {processed}")
    print(f"   Files failed/skipped: {failed}")
    
    return {
        "processed": processed,
        "failed": failed,
    }


if __name__ == "__main__":
    main()
